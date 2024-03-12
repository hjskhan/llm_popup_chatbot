from flask import Flask, render_template, request, session, jsonify, current_app
from flask_session import Session
import os
import openai
from langchain_astradb import AstraDBVectorStore
from langchain_community.llms import OpenAI
from langchain.text_splitter import CharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from dotenv import load_dotenv
from urllib.request import urlopen
from bs4 import BeautifulSoup
import requests
from astrapy.db import AstraDB
import secrets
import random
import string


app = Flask(__name__)
app.config["SESSION_PERMANENT"] = False # Session is not permanent
app.config["SESSION_TYPE"] = "filesystem" # Store session data in file system
app.config["AstraVectorStore"] = None # Store the vector store in session
Session(app)

# Load environment variables from .env file
load_dotenv(override=True)

ASTRA_DB_APPLICATION_TOKEN = os.getenv('ASTRA_DB_APPLICATION_TOKEN')
ASTR_DB_ID = os.getenv('ASTR_DB_ID')
ASTRA_DB_API_ENDPOINT = os.getenv('ASTRA_DB_API_ENDPOINT')
OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
app.secret_key = secrets.token_hex() # Generate a random secret key



####################################################
# Function to generate random string for table_name

def generate_random_string(length=10):
    characters = string.ascii_lowercase + string.digits + '_'
    random_string = ''.join(random.choice(characters) for _ in range(length - 2))
    return random.choice(string.ascii_lowercase) + random_string + random.choice(string.ascii_lowercase)
table_name = generate_random_string()
# table_name = 'astra_vector_demo'
####################################################
# AstraDB connection and vector store initialization
def initialize_astra_vector_store():
    # cassio.init(token=ASTRA_DB_APPLICATION_TOKEN, database_id=ASTR_DB_ID)
    
    embedding = OpenAIEmbeddings(openai_api_key=OPENAI_API_KEY)

    astra_vector_store = AstraDBVectorStore(
    embedding=embedding,
    collection_name=table_name,
    api_endpoint=ASTRA_DB_API_ENDPOINT,
    token=ASTRA_DB_APPLICATION_TOKEN
)
    print('----------------------table created as table_name:', table_name)
    return astra_vector_store
####################################################
####################################################
# Function to preprocess url
def preprocess_url(url):
    texts = []
    print('---------------------- Inside preprocess_url')
    response = requests.get(url)
    soup = BeautifulSoup(response.text, 'html.parser')
    
    # kill all script and style elements
    for script in soup(["script", "style"]):
        script.decompose()    # rip it out

    raw_text = soup.get_text()

    text_splitter = CharacterTextSplitter(
        separator='\n',
        chunk_size=900,
        chunk_overlap=200,
        length_function=len
    )
    texts = text_splitter.split_text(raw_text)
    print('---------------------- chunks created from url')
    print('---------------------- preprocess_url done')
    return texts

# Function to preprocess the uploaded file
def preprocessor_files(uploaded_files):
    texts = []
    if uploaded_files:
        print('---------------------- Inside preprocessor')
        if uploaded_files.filename.endswith('.pdf'):
            texts.extend(preprocess_pdf(uploaded_files))
        elif uploaded_files.filename.endswith(('.doc', '.docx')):
            texts.extend(preprocess_word(uploaded_files))
        elif uploaded_files.filename.endswith('.txt'):
            texts.extend(preprocess_text(uploaded_files))
        elif uploaded_files.filename.endswith('.md'):
            texts.extend(preprocess_markdown(uploaded_files))
        elif uploaded_files.filename.endswith(('.html', '.htm')):
            texts.extend(preprocess_html(uploaded_files))
        elif uploaded_files.filename.endswith(('.pptx')):
            texts.extend(preprocess_pptx(uploaded_files))

    return texts



# Sub-Function to preprocess powerpoint file (.pptx)
def preprocess_pptx(uploaded_file):
    prs = Presentation(uploaded_file)
    raw_text = ''

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                raw_text += shape.text + '\n'    

    text_splitter = CharacterTextSplitter(
        separator='\n',
        chunk_size=900,
        chunk_overlap=200,
        length_function=len
    )
    texts = text_splitter.split_text(raw_text)
    return texts

# Sub-Function to preprocess plain text file (.txt)
def preprocess_text(uploaded_file):
    raw_text = uploaded_file.read().decode('utf-8')

    text_splitter = CharacterTextSplitter(
        separator='\n',
        chunk_size=900,
        chunk_overlap=200,
        length_function=len
    )
    texts = text_splitter.split_text(raw_text)
    return texts

# Sub-Function to preprocess markdown file (.md)
def preprocess_markdown(uploaded_file):
    
    raw_text = uploaded_file.read().decode('utf-8')

    # Additional markdown processing logic if needed

    text_splitter = CharacterTextSplitter(
        separator='\n',
        chunk_size=900,
        chunk_overlap=200,
        length_function=len
    )
    texts = text_splitter.split_text(raw_text)
    return texts

# Sub-Function to preprocess HTML file (.html)
def preprocess_html(uploaded_file):
    # Additional HTML processing logic if needed
    raw_text = uploaded_file.read().decode('utf-8')

    text_splitter = CharacterTextSplitter(
        separator='\n',
        chunk_size=900,
        chunk_overlap=200,
        length_function=len
    )
    texts = text_splitter.split_text(raw_text)
    return texts



# Sub-Function to preprocess pdf
def preprocess_pdf(uploaded_file):
    pdf_reader = PdfReader(uploaded_file)
    raw_text = ''
    for page in pdf_reader.pages:
        content = page.extract_text()
        if content:
            raw_text += content

    text_splitter = CharacterTextSplitter(
        separator='\n',
        chunk_size=900,
        chunk_overlap=200,
        length_function=len
    )
    texts = text_splitter.split_text(raw_text)
    print('---------------------- Inside preprocess_pdf')
    return texts

# Sub-Function to preprocess word document
def preprocess_word(uploaded_file):
    document = Document(uploaded_file)
    raw_text = ''
    for paragraph in document.paragraphs:
        raw_text += paragraph.text

    text_splitter = CharacterTextSplitter(
        separator='\n',
        chunk_size=900,
        chunk_overlap=200,
        length_function=len
    )
    texts = text_splitter.split_text(raw_text)
    return texts
####################################################


####################################################
# Function to perform query from AstraDB
def perform_query(query_text, astra_vector_store):
    vectorDB_answer = astra_vector_store.similarity_search_with_score(query_text, k=1)
    res, score = vectorDB_answer[0]
    return res.page_content, score


# Function to perform query from chatcompletion
def perform_query_chat(message_history):
    response = openai.chat.completions.create(
                model="gpt-3.5-turbo",
                messages = message_history,
                temperature=0.7,
                max_tokens=800,
                top_p=0.95,
                frequency_penalty=0,
                presence_penalty=0,
                stop=None
                # stream=True
                )
    return response
####################################################


message_history = []# message history
# initial message
message_text = [{"role":"system","content":"You are an AI assistant that helps people by answering the questions asked."}]
message_history.extend(message_text)
# Initiate Flask routes for chatbot
@app.route('/')
def index():
    current_app.config["AstraVectorStore"] = initialize_astra_vector_store() # Store the vector store in session
    print('----------------------astraDB initialized')
    
    return render_template('index.html', uploaded_files=session.get('uploaded_files', []))

# Route to upload file
@app.route('/upload', methods=['POST'])
def upload():
    
    if 'file' in request.files:
        file = request.files['file']
        if file.filename != '':
            uploaded_files = session.get('uploaded_files', [])
            uploaded_files.append(file.filename)
            session['uploaded_files'] = uploaded_files
            # preprocess the file and store the texts in session
            session['texts'] = preprocessor_files(file)
            texts = session.get('texts')
            astra_vector_store = current_app.config["AstraVectorStore"]
            astra_vector_store.add_texts(texts)
            # astra_vector_store.add_documents(texts)
            print('----------------------texts added to astraDB')
            current_app.config["AstraVectorStore"] = astra_vector_store
            return 'File uploaded and preprocessed successfully!'
    return 'No file selected!'

# Route to upload URL
@app.route('/upload_url', methods=['POST'])
def upload_url():
    
    if 'url' in request.form:
        url = request.form['url']
        url_links = session.get('url_links', [])
        url_links.append(url)
        session['url_links'] = url_links
        
        print('---------------------- URLs stored in session:', session['url_links'])
        # preprocess the url and store the texts in session
        for i in url_links:
            print('---------------------- url:', i)

        session['texts'] = preprocess_url(url)
        texts = session.get('texts')
        astra_vector_store = current_app.config["AstraVectorStore"]
        astra_vector_store.add_texts(texts)
        print('----------------------texts added to astraDB')
        current_app.config["AstraVectorStore"] = astra_vector_store
        return 'URL uploaded successfully!'
    return 'No URL selected!'


# Chatbot route--main route for chatbot--handels both GET and POST requests
@app.route('/chat', methods=['GET', 'POST'])
def chatbot():
    
    if request.method == 'POST': # if there is file upload then preprocess it\
        user_message = request.form.get('message')

        # Check if the uploaded files are available in session
        if ((session.get('uploaded_files') is not None) or (session.get('url_links') is not None)) and (current_app.config["AstraVectorStore"] is not None):
            astra_vector_store = current_app.config["AstraVectorStore"]
            print('----------------------astra_vector_index is not None')
            
            # Perform initial query on VectorDB
            vectorDB_answer, score = perform_query(user_message, astra_vector_store)
            # vectorDB_answer = astra_vector_store.similarity_search_with_score(user_message, k=1)
            print('----------------------vectorDB query performed')
            # print('----------------------vectorDB_answer:', vectorDB_answer)
                
            if score >= 0.89:
                prompt = f"""
                        Answer the question as detailed as possible from the provided context, make sure to provide all the details, if the answer is not in \
                        provided context just say, "answer is not available in the context", don't provide the wrong answer\n\n \
                        Context:\n {vectorDB_answer}\n \
                        Question: \n{user_message}\n \
                        """
                message_history.append({"role": "user", "content": prompt})
                response = perform_query_chat(message_history)
                
                # append assistant message to message history
                message_history.append({"role":"assistant","content":response.choices[0].message.content})
                return jsonify({"response": response.choices[0].message.content})
            
            else:
                latest_assistant_content = next(item['content'] for item in reversed(message_history) if item['role'] == 'assistant')
                # print('----------------------latest_assistant:',latest_assistant_content)
                contextualize_prompt = f"""Given a content and a user question, if the user question is in reference to the context in\
                  the given content or if it asks for summarization, elaboration, detailed explanation, or related queries related about the content, formulate a\
                    standalone question that can be understood without the given content. Otherwise, return 'NO'. Do NOT answer the question;\
                      just reformulate it if needed.\n\n
                                        user question:\n {user_message}\n \
                                        content: {latest_assistant_content}
                                        """         

                prompt = [{"role":"system","content":"You are an AI assistant that helps people by answering the questions asked."},
                 {"role": "user", "content": contextualize_prompt}]
                # print('----------------------contextualize_prompt:', contextualize_prompt)
                response = perform_query_chat(prompt)

                if response.choices[0].message.content == 'NO':
                    print('----------------------NO')
                    message_history.append({"role":"user","content":user_message})
                    message_history.append({"role":"assistant","content":"No related answer is available in the files or url uploaded!"})
                    return jsonify({"response": "No related answer is available in the files or url uploaded!"})
                else:
                    print('----------------------YES')
                    message_history.append({"role":"user","content":response.choices[0].message.content})
                    response = perform_query_chat(message_history)
                    message_history.append({'role': 'assistant', 'content': response.choices[0].message.content})
                    return jsonify({"response": response.choices[0].message.content})

        else:
            return jsonify({"response": "Please upload a file or URL to start the conversation!"})


if __name__ == '__main__':
    app.run(debug=True)